from pptx import Presentation
import sys
# REGEXXXXXX from re import 

# open presentation from arguments
# iterate over slides
# extract image from slide
# extract presenter notes from slide
# trim presenter notes to essential information
## Title
## Culture / movement
## Geographic origin
## Artist
## Date
## Artistic period
## Location
## Significance

def scrape_presentation(path):
    p = Presentation(path)

    i = 0
    for slide in p.slides:
        #print("----------\n"+slide.notes_slide.notes_text_frame.text)
        if i > 0:
            open('text_test_% s.txt' % i, 'w').write(slide.notes_slide.notes_text_frame.text)
            open('image_test_% s.jpeg' % i, 'wb').write(slide.shapes[0].image.blob)
        i = i + 1

if __name__ == '__main__':
    globals()[sys.argv[1]](sys.argv[2])
